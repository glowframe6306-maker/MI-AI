import html
import json
import os
import re
import shutil
import traceback
import sys
import uuid
import urllib.error
import urllib.parse
import urllib.request
from pathlib import Path
from pprint import pprint


try:
    from google import genai
    print("[BOOT] google-genai imported successfully")
    try:
        # Try to read the installed package version
        try:
            import importlib.metadata as importlib_metadata
            gg_version = importlib_metadata.version("google-genai")
        except Exception:
            try:
                import pkg_resources
                gg_version = pkg_resources.get_distribution("google-genai").version
            except Exception:
                gg_version = "unknown"
        print(f"[BOOT] Python executable: {sys.executable}")
        print(f"[BOOT] google-genai version: {gg_version}")
        print("[BOOT] Gemini SDK loaded successfully")
    except Exception:
        print("[BOOT] Could not determine google-genai version")
except Exception as import_error:
    print(f"[BOOT] Failed to import google-genai: {import_error}")
    import traceback as tb
    tb.print_exc()
    genai = None


def debug_gemini_direct():
    from dotenv import load_dotenv

    load_dotenv(dotenv_path=Path(__file__).resolve().parent / ".env")
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash").strip()

    print("=== Direct Gemini environment ===")
    print("GEMINI_API_KEY length:", len(api_key))
    print("GEMINI_MODEL:", model)

    if not api_key or genai is None:
        print("Gemini client unavailable")
        return

    client = genai.Client(api_key=api_key)
    try:
        response = client.models.generate_content(
            model=model,
            contents="Say hello in one short sentence.",
        )
        print("=== Gemini generate_content ===")
        print(getattr(response, "text", ""))
    except Exception as exc:
        print("Gemini request failed")
        traceback.print_exc()
        print("status_code:", getattr(exc, "status_code", None))
        print("body:", getattr(exc, "body", None))


from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from dotenv import load_dotenv


def _safe_console_text(value) -> str:
    if value is None:
        return ""
    text = str(value)
    encoding = getattr(sys.stdout, "encoding", None) or "utf-8"
    try:
        return text.encode(encoding, errors="replace").decode(encoding, errors="replace")
    except Exception:
        return text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")


def _load_environment() -> list[str]:
    candidates = []
    for base in [Path.cwd(), Path(__file__).resolve().parent, Path(__file__).resolve().parent.parent]:
        candidate = base / ".env"
        candidates.append(candidate)
        if candidate.exists():
            load_dotenv(dotenv_path=candidate, override=False)
    return [str(candidate) for candidate in candidates]


_load_environment()

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

app = Flask(__name__)
CORS(app)

gemini_client = None
gemini_client_api_key = None


def _get_gemini_api_key():
    """Read the Gemini API key from environment variables."""
    _load_environment()
    value = os.getenv("GEMINI_API_KEY", "").strip()
    if value:
        return value
    return ""


# Removed OpenAI compatibility helper to enforce Gemini-only usage


def _get_gemini_model():
    _load_environment()
    model = os.getenv("GEMINI_MODEL", "").strip()
    return model or "gemini-2.5-flash"


def _get_available_models(client):
    """Fetch list of available models that support generateContent."""
    if client is None:
        print("Client is None, cannot fetch models")
        return []

    try:
        print("Fetching available models from Gemini API...")
        models = client.models.list()
        if models is None:
            print("Model list returned None")
            return []

        available = []
        for model in models:
            model_name = getattr(model, 'name', None) or str(model)
            model_id = model_name.replace('models/', '').strip()
            if not model_id:
                continue

            model_lower = model_id.lower()
            if any(tag in model_lower for tag in ['preview', 'tts', 'exp', 'beta', 'alpha']):
                continue

            supported_actions = getattr(model, 'supported_actions', None)
            if supported_actions is None:
                print(f"  Skipping {model_id} because supported_actions is unknown")
                continue

            if 'generateContent' in supported_actions:
                available.append(model_id)
                print(f"  {model_id}")

        print(f"Total models supporting generateContent: {len(available)}")
        return available
    except Exception as e:
        print(f"Error fetching models: {e}")
        traceback.print_exc()
        return []


def _select_best_model(client, requested_model=None):
    """Select the best available model, with priority order."""
    priority_models = ["gemini-2.5-flash"]
    if requested_model:
        requested_model = requested_model.replace('models/', '').strip()
        if requested_model and requested_model not in priority_models:
            priority_models.insert(0, requested_model)

    available_models = _get_available_models(client)

    if not available_models:
        print("Warning: Could not fetch available models from API")
        fallback = [m for m in priority_models if m != requested_model]
        return requested_model or "gemini-2.5-flash", fallback

    print(f"\nAvailable models: {available_models}")
    print(f"Priority order: {priority_models}")

    for model in priority_models:
        if model in available_models:
            fallback = [m for m in available_models if m != model]
            print(f"Selected model: {model}")
            print(f"Fallback models: {fallback}")
            return model, fallback

    selected = available_models[0]
    fallback = available_models[1:]
    print(f"No priority model available, using first available: {selected}")
    print(f"Selected model: {selected}")
    print(f"Fallback models: {fallback}")
    return selected, fallback


def _format_gemini_error_message(exc):
    """Format Gemini exception into user-friendly error message."""
    error_str = str(exc).lower()
    
    # Check for quota/rate limit errors from google.genai.errors.ClientError
    if "429" in str(exc) or "resource_exhausted" in error_str or "quota" in error_str or "rate_limit" in error_str:
        return "The Gemini service is temporarily unavailable because of quota or rate limits. Your free tier quota has been exceeded. Please wait a few hours or upgrade your plan."
    
    # Check for authentication errors
    if "401" in str(exc) or "403" in str(exc) or "invalid" in error_str or "authentication" in error_str:
        return "The Gemini API key is invalid or not authorized. Please verify the key and permissions."
    
    # Check for model not found errors  
    if "404" in str(exc) or "not_found" in error_str or "model_not_found" in error_str:
        return "The requested Gemini model is not available. Please try a supported model."
    
    # Check for server errors
    if "500" in str(exc) or "503" in str(exc) or "server_error" in error_str or "internal_error" in error_str:
        return "The Gemini service is temporarily unavailable. Please try again shortly."
    
    # Fallback to the original exception message if it's not too long
    exc_message = str(exc)
    if len(exc_message) < 200:
        return exc_message
    else:
        # If message is too long, try to extract a meaningful part
        if "RESOURCE_EXHAUSTED" in exc_message:
            return "Quota exceeded on the Gemini API. Please try again later."
        return "An error occurred with the Gemini service. Please try again."


def _get_gemini_client():
    global gemini_client, gemini_client_api_key
    _load_environment()
    api_key = _get_gemini_api_key()
    if not api_key or genai is None:
        gemini_client = None
        gemini_client_api_key = None
        return None

    if gemini_client is not None and gemini_client_api_key == api_key:
        return gemini_client

    try:
        print("=== Gemini environment ===")
        print("GEMINI_API_KEY length:", len(api_key))
        print("GEMINI_MODEL:", _get_gemini_model())
        print("Environment:", os.getcwd())
        gemini_client = genai.Client(api_key=api_key)
        gemini_client_api_key = api_key
        print("=== Gemini client initialized ===")
    except Exception as exc:
        gemini_client = None
        gemini_client_api_key = None
        print("Gemini client initialization failed")
        traceback.print_exc()
        print("Gemini init error:", exc)

    return gemini_client


def _build_gemini_contents(messages_payload):
    contents = []
    print(f"\n=== Building Gemini contents ===")
    print(f"Input messages: {len(messages_payload)}")
    for i, message in enumerate(messages_payload):
        role = (message.get("role") or "").lower()
        content = message.get("content", "") or ""
        print(f"Message {i}: role={role}, content_len={len(content)}")
        if not content:
            print(f"  -> Skipping empty content")
            continue
        if role == "system":
            mapped_role = "user"
            text = f"System instruction: {content}"
            print(f"  -> Mapped system to {mapped_role}")
            contents.append({"role": mapped_role, "parts": [{"text": text}]})
        elif role == "assistant":
            mapped_role = "model"
            text = content
            print(f"  -> Mapped assistant to {mapped_role}")
            contents.append({"role": mapped_role, "parts": [{"text": text}]})
        elif role in {"user", "model"}:
            print(f"  -> Using role as-is: {role}")
            contents.append({"role": role, "parts": [{"text": content}]})
        else:
            mapped_role = "user"
            print(f"  -> Mapped unknown role {role} to {mapped_role}")
            contents.append({"role": mapped_role, "parts": [{"text": content}]})
    print(f"Final contents: {len(contents)} items")
    return contents


def _generate_gemini_content(client, contents, model_name=None, config=None):
    if client is None:
        raise RuntimeError("Gemini client is not available")

    requested_model = (model_name or _get_gemini_model()).strip()
    
    # Get best available model and fallback list
    selected_model, fallback_models = _select_best_model(client, requested_model)
    
    print(f"\n=== Gemini content generation ===")
    print(f"Requested model: {requested_model}")
    print(f"Selected model: {selected_model}")
    print(f"Fallback models: {fallback_models}")
    print(f"Contents: {len(contents)} messages")
    print(f"Config: {config}")
    
    # Use selected model, then fallback to others if it fails
    fallback_models_to_try = [selected_model]
    for model in fallback_models:
        if model != selected_model and model not in fallback_models_to_try:
            fallback_models_to_try.append(model)

    last_error = None
    for candidate_model in fallback_models_to_try:
        try:
            print(f"\n>>> Trying model: {candidate_model}")
            response = client.models.generate_content(
                model=candidate_model,
                contents=contents,
                config=config,
            )
            print(f"<<< Model {candidate_model} succeeded")
            print(f"Response type: {type(response)}")
            print(f"Response attributes: {dir(response)}")
            return response
        except Exception as exc:
            last_error = exc
            status_code = getattr(exc, "status_code", None) or getattr(exc, "code", None)
            body = getattr(exc, "body", None)
            error_message = str(exc)
            
            print(f"!!! Model {candidate_model} failed")
            print(f"Exception type: {type(exc).__name__}")
            print(f"Status code: {status_code}")
            print(f"Error message: {error_message}")
            print(f"Body: {body}")
            
            if isinstance(body, dict):
                error = body.get("error") or body
                code = (error.get("code") if isinstance(error, dict) else None) or ""
                error_type = (error.get("type") if isinstance(error, dict) else None) or ""
                message = (error.get("message") if isinstance(error, dict) else None) or ""
            else:
                code = ""
                error_type = ""
                message = ""
            
            # Check if this is a retryable error
            is_quota_error = (
                status_code in {429, 500, 503} or
                code in {"resource_exhausted", "quota_exceeded", "rate_limit_exceeded", "insufficient_quota"} or
                error_type in {"resource_exhausted", "quota_exceeded", "rate_limit_exceeded", "insufficient_quota"} or
                "quota" in error_message.lower() or
                "rate_limit" in error_message.lower() or
                "resource_exhausted" in error_message.lower()
            )
            
            print(f"Is retryable (quota/rate limit): {is_quota_error}")
            
            # If not retryable, raise immediately
            if not is_quota_error:
                print(f"Non-retryable error, raising")
                raise
            
            # If this is the last model, raise
            if candidate_model == fallback_models_to_try[-1]:
                print(f"Last model in list, raising")
                raise
            
            print(f"Retrying with next model...")

    if last_error is not None:
        raise last_error
    raise RuntimeError("Gemini content generation failed")


SUPABASE_URL = os.getenv("SUPABASE_URL", "").rstrip("/")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
LOCAL_PERSISTENCE_FILE = Path(__file__).with_name("local_persistence.json")
UPLOAD_DIR = Path(__file__).with_name("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)
FRONTEND_DIR = Path(__file__).resolve().parent.parent / "frontend"
FRONTEND_INDEX = FRONTEND_DIR / "index.html"


def _load_local_store():
    if LOCAL_PERSISTENCE_FILE.exists():
        try:
            return json.loads(LOCAL_PERSISTENCE_FILE.read_text(encoding="utf-8"))
        except Exception:
            return {"conversations": [], "messages": []}
    return {"conversations": [], "messages": []}


def _save_local_store(store):
    LOCAL_PERSISTENCE_FILE.write_text(json.dumps(store, indent=2), encoding="utf-8")


def _safe_file_name(name):
    safe = re.sub(r'[^A-Za-z0-9._-]+', '_', (name or 'upload').strip())
    return safe or 'upload'


def _allowed_file_type(filename, content_type):
    allowed_exts = {'.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.csv', '.pptx', '.ppt', '.png', '.jpg', '.jpeg', '.webp'}
    allowed_mimes = {
        'application/pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/msword',
        'text/plain',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.ms-excel',
        'text/csv',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'application/vnd.ms-powerpoint',
        'image/png',
        'image/jpeg',
        'image/webp',
    }
    ext = Path(filename or '').suffix.lower()
    return ext in allowed_exts and (not content_type or content_type in allowed_mimes or ext in {'.png', '.jpg', '.jpeg', '.webp'})


def _needs_live_web_search(text):
    if not text:
        return False
    lowered = text.lower()
    patterns = [
        r"\b(latest|current|today|now|live|breaking|recent|this week|this month|upcoming|tonight|weather|news|stock|price|forecast|schedule|result|score|release|latest news|what['’]?s happening|who won|what time|what day|when does|where is|how much is)\b",
        r"\b(what['’]?s the latest|what['’]?s happening|recent updates|up to date|today['’]?s|current status)\b",
        r"\b(search|find|look up|google|online|internet|official website|official docs|documentation|github|youtube|image|video|pdf|research|map|maps|restaurant|hotel|hospital|bank|mosque|temple|school|product|price|cheapest|compare|best product|buy|flight|hotel|shop|shopping|location|address)\b"
    ]
    return any(re.search(pattern, lowered) for pattern in patterns)


def _classify_search_intent(text):
    lowered = (text or "").lower()
    if any(keyword in lowered for keyword in ["image", "picture", "photo", "images", "show me a picture", "show me", "kaputa"]):
        return "images"
    if any(keyword in lowered for keyword in ["video", "tutorial", "interview", "documentary", "song", "gameplay", "movie trailer", "youtube"]):
        return "videos"
    if any(keyword in lowered for keyword in ["pdf", "research paper", "manual", "documentation", "user guide", "guide", "doc"]):
        return "documents"
    if any(keyword in lowered for keyword in ["price", "buy", "cheapest", "compare", "best product", "product", "shopping", "shop"]):
        return "shopping"
    if any(keyword in lowered for keyword in ["restaurant", "hotel", "hospital", "bank", "police", "mosque", "temple", "school", "nearest", "map", "maps", "location", "address"]):
        return "maps"
    if any(keyword in lowered for keyword in ["news", "latest", "breaking", "current events", "today", "tonight", "recent"]):
        return "news"
    if any(keyword in lowered for keyword in ["github", "npm", "pypi", "firebase", "supabase", "cloudflare", "vercel", "mdn", "stackoverflow", "official docs", "official website"]):
        return "developer"
    return "web"


def _fetch_web_search_results(query, limit=4):
    if not query:
        return []
    encoded_query = urllib.parse.quote(query)
    search_url = f"https://html.duckduckgo.com/html/?q={encoded_query}"
    request = urllib.request.Request(
        search_url,
        headers={"User-Agent": "Mozilla/5.0"},
    )
    try:
        with urllib.request.urlopen(request, timeout=8) as response:
            html_text = response.read().decode("utf-8", "ignore")
    except Exception:
        return []

    results = []
    for match in re.finditer(r'<a rel="nofollow" class="result__a" href="(.*?)"(.*?)>(.*?)</a>', html_text, flags=re.S):
        href = html.unescape(match.group(1) or "")
        title = re.sub(r"<.*?\>", "", match.group(3) or "")
        title = html.unescape(re.sub(r"\s+", " ", title)).strip()
        if href and title:
            results.append({"title": title, "url": href})
        if len(results) >= limit:
            break
    return results


def _build_search_links(query, intent):
    encoded_query = urllib.parse.quote(query)
    if intent == "images":
        return [
            f"https://www.google.com/search?tbm=isch&q={encoded_query}",
            f"https://duckduckgo.com/?q={encoded_query}&iax=images&ia=images",
        ]
    if intent == "videos":
        return [
            f"https://www.youtube.com/results?search_query={encoded_query}",
            f"https://www.google.com/search?tbm=vid&q={encoded_query}",
        ]
    if intent == "maps":
        return [
            f"https://www.google.com/maps/search/{encoded_query}",
            f"https://www.google.com/search?q={encoded_query}+map",
        ]
    if intent == "news":
        return [
            f"https://news.google.com/search?q={encoded_query}",
            f"https://www.google.com/search?q={encoded_query}+news",
        ]
    if intent == "documents":
        return [
            f"https://www.google.com/search?q={encoded_query}+pdf",
            f"https://www.google.com/search?q={encoded_query}+documentation",
        ]
    if intent == "shopping":
        return [
            f"https://www.google.com/search?q={encoded_query}+price+official+store",
            f"https://www.google.com/search?q={encoded_query}+buy",
        ]
    if intent == "developer":
        return [
            f"https://www.google.com/search?q={encoded_query}+official+documentation",
            f"https://www.google.com/search?q={encoded_query}+github",
        ]
    return [
        f"https://www.google.com/search?q={encoded_query}",
        f"https://duckduckgo.com/?q={encoded_query}",
        f"https://en.wikipedia.org/wiki/Special:Search?search={encoded_query}",
    ]


def _build_live_search_context(user_message):
    if not _needs_live_web_search(user_message):
        return ""

    query = user_message.strip()
    intent = _classify_search_intent(query)
    results = _fetch_web_search_results(query, limit=4)
    search_links = _build_search_links(query, intent)

    context_lines = ["Live web context:"]
    if results:
        context_lines.append("Web results:")
        for index, item in enumerate(results, 1):
            context_lines.append(f"{index}. {item['title']} - {item['url']}")

    if search_links:
        context_lines.append("Suggested sources:")
        for index, link in enumerate(search_links, 1):
            context_lines.append(f"{index}. {link}")

    return "\n".join(context_lines)


# ==================== REAL-TIME ACCESS SYSTEM ====================

def _needs_real_time_data(text):
    """Detect if user query requires real-time data"""
    if not text:
        return False
    lowered = text.lower()
    
    real_time_patterns = [
        # Time and date
        r"\b(what time|what['']?s the time|current time|what is|what['']?s today|today['']?s date|current date|what day)\b",
        # Weather
        r"\b(weather|temperature|rain|raining|forecast|humidity|wind speed|sunrise|sunset|air quality|uv index|storm|cold|hot|snow|cloudy|sunny)\b",
        # Sports
        r"\b(score|match|game|cricket|football|basketball|tennis|volleyball|formula 1|f1|live|sports|league|table|standings|fixtures|results|player stats|tournament)\b",
        # News and current events
        r"\b(news|breaking|latest|current events|elections|announcements|trending)\b",
        # Crypto and market
        r"\b(bitcoin|ethereum|crypto|price|stock|market|gold|silver|dollar|exchange rate|usd|gbp|eur|inr|lkr)\b",
        # Traffic and transport
        r"\b(traffic|flight status|train|bus|public transport|fuel price|petrol|diesel|current)\b",
        # Disasters and alerts
        r"\b(earthquake|weather warning|alert|emergency|disaster)\b",
        # Time zones and world clock
        r"\b(time zone|world clock|time in|country time|city time)\b",
    ]
    return any(re.search(pattern, lowered) for pattern in real_time_patterns)


def _detect_real_time_type(text):
    """Identify what type of real-time data is needed"""
    lowered = text.lower()

    if any(kw in lowered for kw in ["what time", "current time", "what['']?s the time", "time now", "what time is it"]):
        return "time"
    if any(kw in lowered for kw in ["weather", "temperature", "rain", "raining", "forecast", "humidity", "wind", "sunrise", "sunset", "air quality", "uv", "sunny", "cloudy"]):
        return "weather"
    if any(kw in lowered for kw in ["bitcoin", "ethereum", "crypto", "btc", "eth"]):
        return "crypto"
    if any(kw in lowered for kw in ["stock", "market", "nasdaq", "s&p", "dow"]):
        return "stocks"
    if any(kw in lowered for kw in ["cricket", "football", "soccer", "basketball", "tennis", "volleyball", "f1", "formula 1", "live score", "match score"]):
        return "sports"
    if any(kw in lowered for kw in ["exchange rate", "currency", "usd", "eur", "gbp", "inr", "lkr"]):
        return "currency"
    if any(kw in lowered for kw in ["gold", "silver", "precious", "commodity", "gold price", "silver price"]):
        return "commodities"
    if any(kw in lowered for kw in ["news", "breaking", "latest", "current events", "elections", "announcements", "trending"]):
        return "news"
    if any(kw in lowered for kw in ["flight", "train", "bus", "transport", "fuel price", "petrol", "diesel"]):
        return "transport"
    if any(kw in lowered for kw in ["traffic", "road", "congestion"]):
        return "traffic"
    if any(kw in lowered for kw in ["earthquake", "weather warning", "alert", "emergency", "disaster"]):
        return "alerts"

    return "general"


def _extract_location(text):
    """Extract a likely location from a user query for weather or local queries."""
    if not text:
        return ""
    lowered = text.lower().strip()
    if not lowered:
        return ""

    prefixes = ["in ", "at ", "for ", "near ", "around "]
    for prefix in prefixes:
        pos = lowered.find(prefix)
        if pos != -1:
            candidate = lowered[pos + len(prefix):].strip()
            candidate = re.split(r"\b(today|now|tomorrow|this week|this month|weather|temperature|rain|forecast|time|date)\b", candidate)[0].strip()
            candidate = re.sub(r"[^a-zA-Z0-9\s,.-]", "", candidate).strip()
            if candidate:
                return candidate.title()

    if " in " in lowered:
        parts = lowered.split(" in ", 1)[1]
        parts = re.split(r"\b(today|now|tomorrow|this week|this month|weather|temperature|rain|forecast|time|date)\b", parts)[0].strip()
        parts = re.sub(r"[^a-zA-Z0-9\s,.-]", "", parts).strip()
        if parts:
            return parts.title()

    return ""


def _fetch_current_time():
    """Fetch current time and date"""
    try:
        import datetime
        now = datetime.datetime.now()
        return {
            "time": now.strftime("%H:%M:%S"),
            "date": now.strftime("%A, %B %d, %Y"),
            "timestamp": now.isoformat(),
            "timezone": "Local"
        }
    except Exception as e:
        return {"error": f"Could not fetch time: {str(e)}"}


def _fetch_weather_data(location="auto"):
    """Fetch weather data using Open-Meteo (free, no API key required)."""
    try:
        if location and location != "auto":
            loc_query = urllib.parse.quote(location)
            geocode_url = f"https://geocoding-api.open-meteo.com/v1/search?name={loc_query}&count=1&language=en&format=json"
            request = urllib.request.Request(
                geocode_url,
                headers={"User-Agent": "Mozilla/5.0"}
            )
            with urllib.request.urlopen(request, timeout=5) as response:
                geo_data = json.loads(response.read().decode("utf-8"))
                results = geo_data.get("results") or []
                if results:
                    result = results[0]
                    latitude = result.get("latitude")
                    longitude = result.get("longitude")
                    location_name = result.get("name") or location
                else:
                    raise ValueError("Location not found")
        else:
            latitude = 40.7128
            longitude = -74.0060
            location_name = "New York"

        weather_url = f"https://api.open-meteo.com/v1/forecast?latitude={latitude}&longitude={longitude}&current=temperature_2m,relative_humidity_2m,weather_code,wind_speed_10m&daily=weather_code,temperature_2m_max,temperature_2m_min,precipitation_sum,sunrise,sunset&temperature_unit=fahrenheit&wind_speed_unit=mph&timezone=auto"

        request = urllib.request.Request(
            weather_url,
            headers={"User-Agent": "Mozilla/5.0"}
        )

        with urllib.request.urlopen(request, timeout=5) as response:
            data = json.loads(response.read().decode("utf-8"))
            current = data.get("current", {})
            daily = (data.get("daily") or {})
            return {
                "location": location_name,
                "temperature": f"{current.get('temperature_2m')}°F",
                "humidity": f"{current.get('relative_humidity_2m')}%",
                "wind_speed": f"{current.get('wind_speed_10m')} mph",
                "weather_code": current.get("weather_code"),
                "sunrise": daily.get("sunrise", [None])[0],
                "sunset": daily.get("sunset", [None])[0],
                "source": "Open-Meteo"
            }
    except Exception as e:
        return {"error": f"Could not fetch weather: {str(e)}"}


def _fetch_crypto_prices():
    """Fetch cryptocurrency prices from CoinGecko (free API)"""
    try:
        crypto_url = "https://api.coingecko.com/api/v3/simple/price?ids=bitcoin,ethereum,cardano&vs_currencies=usd&include_market_cap=true&include_24hr_change=true"
        
        request = urllib.request.Request(
            crypto_url,
            headers={"User-Agent": "Mozilla/5.0"}
        )
        
        with urllib.request.urlopen(request, timeout=5) as response:
            data = json.loads(response.read().decode("utf-8"))
            return {
                "bitcoin": data.get("bitcoin", {}).get("usd"),
                "bitcoin_change": data.get("bitcoin", {}).get("usd_24h_change"),
                "ethereum": data.get("ethereum", {}).get("usd"),
                "ethereum_change": data.get("ethereum", {}).get("usd_24h_change"),
                "cardano": data.get("cardano", {}).get("usd"),
                "cardano_change": data.get("cardano", {}).get("usd_24h_change"),
                "source": "CoinGecko"
            }
    except Exception as e:
        return {"error": f"Could not fetch crypto prices: {str(e)}"}


def _fetch_exchange_rates():
    """Fetch currency exchange rates from exchangerate-api (limited free tier)"""
    try:
        # Using a free endpoint that doesn't require auth
        rates_url = "https://api.exchangerate-api.com/v4/latest/USD"
        
        request = urllib.request.Request(
            rates_url,
            headers={"User-Agent": "Mozilla/5.0"}
        )
        
        with urllib.request.urlopen(request, timeout=5) as response:
            data = json.loads(response.read().decode("utf-8"))
            rates = data.get("rates", {})
            
            return {
                "base": "USD",
                "eur": rates.get("EUR"),
                "gbp": rates.get("GBP"),
                "inr": rates.get("INR"),
                "lkr": rates.get("LKR"),
                "jpy": rates.get("JPY"),
                "aud": rates.get("AUD"),
                "timestamp": data.get("time_last_updated"),
                "source": "ExchangeRate-API"
            }
    except Exception as e:
        return {"error": f"Could not fetch exchange rates: {str(e)}"}


def _build_real_time_context(user_message):
    """Build real-time data context based on user query."""
    if not _needs_real_time_data(user_message):
        return ""

    data_type = _detect_real_time_type(user_message)
    location = _extract_location(user_message)
    context_lines = ["Real-Time Data Context:"]

    try:
        if data_type == "time":
            time_data = _fetch_current_time()
            if "error" not in time_data:
                context_lines.append(f"Current time: {time_data['time']}")
                context_lines.append(f"Current date: {time_data['date']}")
                context_lines.append(f"Timezone: {time_data['timezone']}")

        elif data_type == "weather":
            weather_data = _fetch_weather_data(location or None)
            if "error" not in weather_data:
                context_lines.append(f"Location: {weather_data['location']}")
                context_lines.append(f"Temperature: {weather_data['temperature']}")
                context_lines.append(f"Humidity: {weather_data['humidity']}")
                context_lines.append(f"Wind Speed: {weather_data['wind_speed']}")
                if weather_data.get("sunrise"):
                    context_lines.append(f"Sunrise: {weather_data['sunrise']}")
                if weather_data.get("sunset"):
                    context_lines.append(f"Sunset: {weather_data['sunset']}")
                context_lines.append(f"Source: {weather_data['source']}")

        elif data_type == "crypto":
            crypto_data = _fetch_crypto_prices()
            if "error" not in crypto_data:
                context_lines.append(f"Bitcoin: ${crypto_data['bitcoin']} ({crypto_data['bitcoin_change']:+.2f}%)")
                context_lines.append(f"Ethereum: ${crypto_data['ethereum']} ({crypto_data['ethereum_change']:+.2f}%)")
                context_lines.append(f"Cardano: ${crypto_data['cardano']} ({crypto_data['cardano_change']:+.2f}%)")
                context_lines.append(f"Source: {crypto_data['source']}")

        elif data_type == "currency":
            rates_data = _fetch_exchange_rates()
            if "error" not in rates_data:
                context_lines.append(f"Base Currency: {rates_data['base']}")
                context_lines.append(f"EUR: {rates_data['eur']}")
                context_lines.append(f"GBP: {rates_data['gbp']}")
                context_lines.append(f"INR: {rates_data['inr']}")
                context_lines.append(f"LKR: {rates_data['lkr']}")
                context_lines.append(f"Source: {rates_data['source']}")

        elif data_type == "commodities":
            context_lines.append("Commodity price lookup requested.")
            context_lines.append("Use official market or commodity sources where available.")

        else:
            search_results = _fetch_web_search_results(user_message, limit=3)
            if search_results:
                context_lines.append("Latest Results:")
                for idx, result in enumerate(search_results, 1):
                    context_lines.append(f"{idx}. {result['title']} - {result['url']}")

    except Exception:
        context_lines.append("Live data is temporarily unavailable. Please try again shortly.")

    if len(context_lines) > 1:
        return "\n".join(context_lines)
    return ""


def _extract_text_from_file(file_path, filename):
    if not file_path.exists():
        return ''
    if file_path.stat().st_size <= 0:
        return ''

    ext = Path(filename or '').suffix.lower()
    if ext == '.pdf':
        if PyPDF2 is None:
            return ''
        try:
            reader = PyPDF2.PdfReader(file_path)
            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text() or '')
            return '\n\n'.join(text_parts)
        except Exception:
            return ''
    if ext in {'.docx', '.doc'}:
        if DocxDocument is None:
            return ''
        try:
            document = DocxDocument(file_path)
            return '\n'.join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        except Exception:
            return ''
    if ext == '.txt':
        try:
            return file_path.read_text(encoding='utf-8', errors='ignore')
        except Exception:
            return ''
    if ext == '.csv':
        try:
            if pd is not None:
                df = pd.read_csv(file_path)
                return df.to_string(index=False)
            return file_path.read_text(encoding='utf-8', errors='ignore')
        except Exception:
            return ''
    if ext in {'.xlsx', '.xls'}:
        try:
            if pd is not None:
                df = pd.read_excel(file_path)
                return df.to_string(index=False)
            return ''
        except Exception:
            return ''
    if ext in {'.pptx', '.ppt'}:
        if Presentation is None:
            return ''
        try:
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)
            return '\n'.join(text_parts)
        except Exception:
            return ''
    if ext in {'.png', '.jpg', '.jpeg', '.webp'}:
        if Image is None or pytesseract is None:
            return ''
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        except Exception:
            return ''
    return ''


def _get_session_id():
    payload = request.get_json(silent=True) or {}
    if isinstance(payload, dict):
        session_id = payload.get("session_id") or request.headers.get("X-Session-Id") or request.headers.get("X-User-Id")
        if session_id:
            return str(session_id)
    return request.args.get("session_id") or request.headers.get("X-Session-Id") or request.headers.get("X-User-Id") or "anonymous"


def _get_user_email():
    payload = request.get_json(silent=True) or {}
    if isinstance(payload, dict) and payload.get("user_email"):
        return str(payload.get("user_email"))
    return request.headers.get("X-User-Email") or request.headers.get("X-Email") or request.args.get("user_email") or ""


def _get_user_id():
    payload = request.get_json(silent=True) or {}
    if isinstance(payload, dict) and payload.get("user_id"):
        return str(payload.get("user_id"))
    return request.headers.get("X-User-Id") or request.args.get("user_id") or ""


def _get_user_scope():
    return {
        "user_id": _get_user_id(),
        "user_email": _get_user_email(),
        "session_id": _get_session_id(),
    }


def _get_chat_debug_context():
    return {
        "gemini_key_present": bool(_get_gemini_api_key()),
        "gemini_model": _get_gemini_model(),
        "provider": "gemini",
        "chat_route": "/chat",
    }


def _request_supabase(method, path, payload=None, params=None):
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
        raise RuntimeError("SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY must be configured")

    url = f"{SUPABASE_URL}/rest/v1{path}"
    if params:
        separator = '&' if '?' in url else '?'
        query = urllib.parse.urlencode(params, doseq=True)
        url = f"{url}{separator}{query}"

    headers = {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "Content-Type": "application/json",
        "Prefer": "return=representation",
    }

    data = None
    if payload is not None:
        data = json.dumps(payload).encode("utf-8")

    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req, timeout=10) as response:
            body = response.read().decode("utf-8")
            if not body:
                return []
            try:
                return json.loads(body)
            except json.JSONDecodeError:
                return body
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8")
        raise RuntimeError(f"Supabase request failed: {exc.code} {body}") from exc


@app.route("/")
def home():
    if FRONTEND_INDEX.exists():
        return send_from_directory(str(FRONTEND_DIR), "index.html")
    return "MI AI Running 🚀"


@app.route("/debug/chat", methods=["GET"])
def debug_chat():
    return jsonify({
        "message": "python-backend",
        "provider": "gemini",
        "chat_route": "/chat",
        "gemini_key_present": bool(_get_gemini_api_key()),
        "gemini_model": _get_gemini_model(),
    })


@app.route("/api/conversations", methods=["GET", "POST"])
def conversations():
    session_id = _get_session_id()
    scope = _get_user_scope()

    if request.method == "GET":
        try:
            params = {"order": "created_at.desc"}
            if scope.get("user_id"):
                params["user_id"] = f"eq.{scope['user_id']}"
            elif session_id and session_id != "anonymous":
                params["session_id"] = f"eq.{session_id}"
            rows = _request_supabase("GET", "/conversations", params=params)
            conversations_list = []
            for row in rows or []:
                conversations_list.append({
                    "id": row.get("id"),
                    "title": row.get("title", "New chat"),
                    "created_at": row.get("created_at"),
                    "updated_at": row.get("updated_at") or row.get("created_at"),
                    "message_count": row.get("message_count") or 0,
                    "last_preview": row.get("last_preview") or "",
                })
            return jsonify({"conversations": conversations_list})
        except Exception as exc:
            store = _load_local_store()
            return jsonify({"conversations": store.get("conversations", []), "warning": str(exc)})

    payload = request.get_json(silent=True) or {}
    title = payload.get("title", "New chat")
    session_id = payload.get("session_id") or session_id or "anonymous"
    user_email = scope.get("user_email") or ""
    user_id = scope.get("user_id") or session_id

    try:
        row = _request_supabase("POST", "/conversations", payload={
            "title": title,
            "session_id": session_id,
            "user_id": user_id,
            "user_email": user_email,
        })
        if isinstance(row, list) and row:
            row = row[0]
        return jsonify({"conversation": row})
    except Exception as exc:
        store = _load_local_store()
        conversation = {
            "id": f"chat_{len(store.get('conversations', [])) + 1}",
            "title": title,
            "session_id": session_id,
            "user_id": user_id,
            "user_email": user_email,
        }
        store.setdefault("conversations", []).append(conversation)
        _save_local_store(store)
        return jsonify({"conversation": conversation, "warning": str(exc)})


@app.route("/api/conversations/<conversation_id>", methods=["PATCH", "DELETE"])
def conversation_detail(conversation_id):
    scope = _get_user_scope()
    user_id = scope.get("user_id")
    session_id = scope.get("session_id")

    if not user_id and session_id == "anonymous":
        return jsonify({"error": "Authentication required"}), 401

    try:
        existing_rows = _request_supabase("GET", f"/conversations?id=eq.{conversation_id}", params={"select": "id,user_id,session_id"})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500

    existing = (existing_rows or [{}])[0]
    if not existing:
        return jsonify({"error": "Conversation not found"}), 404

    if user_id and existing.get("user_id") and existing.get("user_id") != user_id:
        return jsonify({"error": "Forbidden"}), 403

    if request.method == "PATCH":
        payload = request.get_json(silent=True) or {}
        patch_payload = {}

        if "title" in payload:
            patch_payload["title"] = payload.get("title")
        if "updated_at" in payload:
            patch_payload["updated_at"] = payload.get("updated_at")
        if "last_preview" in payload:
            patch_payload["last_preview"] = payload.get("last_preview")
        if "message_count" in payload:
            patch_payload["message_count"] = payload.get("message_count")

        if not patch_payload:
            return jsonify({"error": "No valid fields provided"}), 400

        try:
            updated = _request_supabase("PATCH", f"/conversations?id=eq.{conversation_id}", payload=patch_payload)
            return jsonify({"conversation": updated})
        except Exception as exc:
            return jsonify({"error": str(exc)}), 500

    try:
        _request_supabase("DELETE", f"/conversations?id=eq.{conversation_id}")
        return jsonify({"deleted": True})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


@app.route("/uploads/<path:filename>")
def serve_upload(filename):
    safe_name = Path(filename).name
    file_path = UPLOAD_DIR / safe_name
    if not file_path.exists():
        return jsonify({"error": "File not found"}), 404
    return send_from_directory(str(UPLOAD_DIR), safe_name, as_attachment=False)


@app.route("/api/upload", methods=["POST"])
def upload_file():
    scope = _get_user_scope()
    user_id = scope.get("user_id") or scope.get("session_id") or "anonymous"
    user_email = scope.get("user_email") or ""
    conversation_id = request.form.get("conversation_id")
    message_id = request.form.get("message_id")

    if not conversation_id:
        return jsonify({"error": "conversation_id is required"}), 400

    uploaded = request.files.get("file")
    if not uploaded:
        return jsonify({"error": "file is required"}), 400

    filename = uploaded.filename or "upload"
    if not filename.strip():
        return jsonify({"error": "file name is invalid"}), 400

    if not _allowed_file_type(filename, uploaded.mimetype):
        return jsonify({"error": "Unsupported file type"}), 400

    try:
        uploaded.stream.seek(0, os.SEEK_END)
        size = uploaded.stream.tell()
        uploaded.stream.seek(0)
    except Exception:
        size = 0

    if size <= 0:
        return jsonify({"error": "File is empty"}), 400

    if size > 20 * 1024 * 1024:
        return jsonify({"error": "File too large"}), 400

    safe_name = _safe_file_name(filename)
    unique_name = f"{uuid.uuid4().hex}_{safe_name}"
    storage_path = UPLOAD_DIR / unique_name

    try:
        uploaded.save(storage_path)
    except Exception:
        return jsonify({"error": "Unable to save uploaded file"}), 400

    try:
        extracted_text = _extract_text_from_file(storage_path, filename)
    except Exception:
        extracted_text = ''

    attachment = {
        "id": f"att_{uuid.uuid4().hex[:8]}",
        "conversation_id": conversation_id,
        "message_id": message_id or None,
        "user_id": user_id,
        "user_email": user_email,
        "file_name": filename,
        "stored_name": unique_name,
        "file_size": storage_path.stat().st_size,
        "file_type": uploaded.mimetype or Path(filename).suffix.lower(),
        "upload_time": __import__('datetime').datetime.utcnow().isoformat() + 'Z',
        "extracted_text": extracted_text,
        "storage_path": str(storage_path),
        "uploaded_file_url": f"/uploads/{unique_name}",
    }

    store = _load_local_store()
    store.setdefault("attachments", []).append(attachment)
    _save_local_store(store)

    return jsonify({"attachment": attachment})


@app.route("/api/conversations/<conversation_id>/attachments", methods=["GET"])
def conversation_attachments(conversation_id):
    store = _load_local_store()
    attachments = [attachment for attachment in store.get("attachments", []) if attachment.get("conversation_id") == conversation_id]
    return jsonify({"attachments": attachments})


@app.route("/api/messages", methods=["GET", "POST"])
def messages():
    session_id = _get_session_id()
    scope = _get_user_scope()

    if request.method == "GET":
        conversation_id = request.args.get("conversation_id")
        if not conversation_id:
            return jsonify({"messages": []})
        try:
            params = {"conversation_id": f"eq.{conversation_id}", "order": "created_at.asc"}
            if scope.get("user_id"):
                params["user_id"] = f"eq.{scope['user_id']}"
            elif session_id and session_id != "anonymous":
                params["session_id"] = f"eq.{session_id}"
            rows = _request_supabase("GET", "/messages", params=params)
            items = []
            for row in rows or []:
                items.append({
                    "id": row.get("id"),
                    "conversation_id": row.get("conversation_id"),
                    "role": row.get("role"),
                    "content": row.get("content"),
                    "attachment_id": row.get("attachment_id"),
                    "attachment_meta": row.get("attachment_meta"),
                    "created_at": row.get("created_at"),
                })
            return jsonify({"messages": items})
        except Exception as exc:
            store = _load_local_store()
            items = [item for item in store.get("messages", []) if item.get("conversation_id") == conversation_id]
            return jsonify({"messages": items, "warning": str(exc)})

    payload = request.get_json(silent=True) or {}
    conversation_id = payload.get("conversation_id")
    content = payload.get("content", "")
    role = payload.get("role", "ai")
    session_id = payload.get("session_id") or session_id or "anonymous"
    user_email = scope.get("user_email") or ""
    user_id = scope.get("user_id") or session_id
    attachment_id = payload.get("attachment_id")
    attachment_meta = payload.get("attachment_meta")

    if not conversation_id:
        return jsonify({"error": "conversation_id is required"}), 400

    try:
        message_payload = {
            "conversation_id": conversation_id,
            "session_id": session_id,
            "user_id": user_id,
            "user_email": user_email,
            "role": role,
            "content": content,
        }
        if attachment_id is not None:
            message_payload["attachment_id"] = attachment_id
        if attachment_meta is not None:
            message_payload["attachment_meta"] = attachment_meta
        row = _request_supabase("POST", "/messages", payload=message_payload)
        if isinstance(row, list) and row:
            row = row[0]
        return jsonify({"message": row})
    except Exception as exc:
        store = _load_local_store()
        message = {
            "id": f"msg_{len(store.get('messages', [])) + 1}",
            "conversation_id": conversation_id,
            "session_id": session_id,
            "user_id": user_id,
            "user_email": user_email,
            "role": role,
            "content": content,
        }
        if attachment_id is not None:
            message["attachment_id"] = attachment_id
        if attachment_meta is not None:
            message["attachment_meta"] = attachment_meta
        store.setdefault("messages", []).append(message)
        _save_local_store(store)
        return jsonify({"message": message, "warning": str(exc)})


@app.route("/api/messages/<message_id>", methods=["PATCH"])
def message_detail(message_id):
    payload = request.get_json(silent=True) or {}
    patch_payload = {}
    if "attachment_id" in payload:
        patch_payload["attachment_id"] = payload.get("attachment_id")
    if "attachment_meta" in payload:
        patch_payload["attachment_meta"] = payload.get("attachment_meta")
    if not patch_payload:
        return jsonify({"error": "No valid fields provided"}), 400

    try:
        updated = _request_supabase("PATCH", f"/messages?id=eq.{message_id}", payload=patch_payload)
        if isinstance(updated, list) and updated:
            updated = updated[0]
        return jsonify({"message": updated})
    except Exception as exc:
        store = _load_local_store()
        for message in store.get("messages", []):
            if message.get("id") == message_id:
                if "attachment_id" in payload:
                    message["attachment_id"] = payload.get("attachment_id")
                if "attachment_meta" in payload:
                    message["attachment_meta"] = payload.get("attachment_meta")
                break
        _save_local_store(store)
        return jsonify({"message": {"id": message_id, **patch_payload}, "warning": str(exc)})



@app.route("/chat", methods=["POST"])
def chat():

    try:

        data = request.get_json()
        print("=== Request payload ===")
        print(data)
        user_message = data.get("message","")
        history = data.get("history") or []
        attachment_ids = data.get("attachment_ids") or []
        attachment_id = data.get("attachment_id")

        if not user_message:
            return jsonify({
                "reply":"Please type a message."
            })


        attachment_context = ""
        attachment_refs = []
        if attachment_ids:
            attachment_refs = [item for item in attachment_ids if item]
        elif attachment_id:
            attachment_refs = [attachment_id]

        if attachment_refs:
            store = _load_local_store()
            matched = []
            for attachment in store.get("attachments", []):
                if attachment.get("id") in attachment_refs:
                    matched.append(attachment)
            if matched:
                attachment_context_parts = []
                for attachment in matched:
                    attachment_context_parts.append(
                        f"\n\nAttached file: {attachment.get('file_name')}\nFile type: {attachment.get('file_type')}\nStorage path: {attachment.get('storage_path') or ''}\nExtracted content:\n{attachment.get('extracted_text') or ''}"
                    )
                attachment_context = "".join(attachment_context_parts)

        messages_payload = [
            {
                "role": "system",
                "content": """

You are MI AI.

Creator:
M.I. Muhammadh

Age of creater: 17 years old

Ambition of creater: Derector of Flight Operations at SpaceX


IMPORTANT:
1. Answer the user's question correctly.
2. Do not invent facts.
3. If you don't know something, say you don't know.
4. Think before answering.
5. Give useful explanations.
6. Be fast and direct.
7. mi ai is an AI assistant created by M.I. Muhammadh.
8. MI AI is must analize the user's question and give the best possible answer.
9. The email of MI AI customer support is miai.customerservice@gmail.com
10. If the user asks about an uploaded file, answer using the attached file content as the primary source of truth.
11. For images, documents, spreadsheets, slides, and PDFs, analyze the complete content and answer any question from it.
12. For multiple files, use the relevant one(s) that match the user's question.
13. If the user asks for current or online information, treat the live web context as a primary source and prefer official or reputable sources when available.
14. When live web context is provided, answer using it to support the response, but clearly say if something could not be verified.
15. If the user asks for images, videos, official websites, documentation, software downloads, GitHub repositories, maps, PDFs, manuals, or guides, include relevant clickable links whenever possible and prefer official sources over third-party websites.
16. For image requests, provide a Google Images or trusted image-search link when appropriate.
17. For video requests, provide a YouTube or trusted video-search link when appropriate.
18. For documentation or software requests, provide the official documentation or download page when possible.
19. For location or map requests, provide a Google Maps link when appropriate.
20. For GitHub or developer-resource requests, provide the official repository or documentation link when possible.
21. Never invent or guess a link; if no reliable source is known, say so clearly.
22. When real-time data is provided (time, weather, sports scores, cryptocurrency prices, exchange rates, etc.), ALWAYS use this data as the source of truth for current information.
23. If the user asks for current time, weather, live sports scores, cryptocurrency prices, currency exchange rates, flight status, news, or any time-sensitive information, prioritize the real-time data provided.
24. Never provide outdated or cached information when real-time data is available.
25. If real-time data is temporarily unavailable, clearly inform the user instead of providing potentially false information.
26. For weather queries, provide temperature, humidity, wind speed, and any relevant forecasts when available.
27. For sports queries, provide live scores, match statistics, and league standings when available.
28. For cryptocurrency queries, provide current prices and 24-hour change percentages.
29. For currency exchange queries, provide current conversion rates between major currencies.
30. Always mention the source and timestamp of real-time data when appropriate.
31. By default, always provide a complete, detailed, and helpful answer. Do not intentionally be brief unless the user explicitly asks for a short response.
32. If the user does not ask for a short answer, explain the topic clearly, include important details, examples when helpful, steps when explaining a process, and warnings or notes when relevant.
33. For explanations, tutorials, coding questions, programming, comparisons, troubleshooting, research, science, mathematics, history, travel, and product recommendations, give a fuller and more comprehensive answer.
34. Only use a short answer if the user explicitly requests one, for example with words like: short answer, briefly, one sentence, summarize, TL;DR, or concise.
35. Do not stop after only one or two sentences when more useful information can be provided.
LANGUAGE RULE:
- Detect the language of the user's latest message.
- Reply ONLY in that language.
- English message = English reply only.
- Sinhala message = Sinhala reply only.
- Never mix languages unless the user mixes first.
- Always reply in the same script as the user's message.
- Always Use 100% correct words and grammar in replies.
- If the user writes in Sinhala letters, reply using Sinhala letters (සිංහල අකුරු).
- If the user writes in Tamil letters, reply using Tamil letters (தமிழ் எழுத்துக்கள்).
- If the user writes in English, reply using English.       
- If user uses another language, reply in that language.
- Do not translate unless asked.    
- if user writes in Singlish (Roman Sinhala), reply using Sinhala letters (සිංහල අකුරු).
- Do not use Singlish when user writes Sinhala.
- Do not use Tanglish when user writes Tamil.
- Never mix languages unless the user mixes them first.
- If user ask any question in any language, MI AI must reply in the same language and script as the user's question.
- ALWAYS follow the above language rules.
You can help with:
- Coding
- Science
- Maths
- Technology
- General knowledge
- Explanations
- Writing
- speech
- Exam preparations
- Learning new topics
- Language translations
- Learning new languages
- Learning new skills
- Life advice
- Learning new hobbies
- Learning new things
- Learning new subjects
- Learning new technologies
- Learning new programming languages
- Learning new frameworks
- Learning new tools
- Basic to advanced level topics
- Creating content
- Debugging code
- Giving step by step solutions
- Giving detailed explanations
- Giving concise answers
- Giving simple answers
- Giving easy to understand answers
- Giving in depth answers
- Giving short answers
- Giving long answers
- Giving examples
- Giving code examples
- Giving real life examples
- Giving practical examples
- Giving theoretical examples
- Giving mathematical examples
- Giving scientific examples
- Giving historical examples
- Giving philosophical examples
- Giving detailed explanations with examples
- Giving concise explanations with examples
- Giving simple explanations with examples
- Genarating new images based on user prompts
- Genarating new text based on user prompts
- Genarating new code based on user prompts
- Genarating new content based on user prompts
- Gebarate image captions based on user prompts
- Genarating new ideas based on user prompts
- Genarating new concepts based on user prompts
- Genarating new solutions based on user prompts
- and much more.

Your style:
Helpful, really smart and friendly.
You are MI AI.

IMPORTANT:
1. Answer correctly.
2. Do not invent facts.
3. If you don't know, say you don't know.
4. Be fast and direct.
5. Do not mention MI AI in replies.
6. Your name is MI AI
7. Your creater is M.I. Muhammadh
8. Always Must give full and complete answers to the user's question.
9. Use emojis when useful and appropriate.
10. Always use emojis in end of your answers when useful and appropriate.dont use in middle of sentences.
11. You must use only one emojy in each sentence and only at the end of the sentence when useful and appropriate.
12. Always follow the above rules and instructions.

LANGUAGE RULE:
- Detect the user's language.
- Reply ONLY in that language.
- Sinhala message = Sinhala reply.
- English message = English reply.
- If user uses another language, reply in that language.
- Do not translate unless asked.
- Detect the user's language.
- Reply ONLY in the same language and script.
- Sinhala typed in Sinhala letters = reply using Sinhala letters (සිංහල අකුරු).
- Tamil typed in Tamil letters = reply using Tamil letters (தமிழ் எழுத்துக்கள்).
- English typed in English = reply using English.
- Do not use Singlish when user writes Sinhala.
- Do not use Tanglish when user writes Tamil.
- Never mix languages unless the user mixes them first.
LANGUAGE RULE:
- Detect the exact writing style of the user's message.
- Always reply using the same language AND same script.

Examples:
- User: "ඔයා කොහොමද?"
  Reply: "මම හොඳින් ඉන්නවා."

- User: "oya kohomada?"
  Reply: "mama hondin innawa."

- User: "How are you?"
  Reply: "I am doing well."

- User: "நீ எப்படி இருக்கிறாய்?"
  Reply: "நான் நன்றாக இருக்கிறேன்."
- Sinhala letters input = Sinhala letters output only.
- Singlish input = Singlish output only.
- English input = English output only.
- Tamil letters input = Tamil letters output only.
- Do not convert Sinhala letters into Singlish.
- Do not convert Singlish into Sinhala letters.
- Do not mix scripts.
LANGUAGE RULE:

- Detect the user's language.
- If the user message is Sinhala OR Singlish (Roman Sinhala),
  always reply using Sinhala Unicode letters.

Examples:

User: "mata udaw karanna"
Reply: "මම උදව් කරන්නම්."

User: "මට උදව් කරන්න"
Reply: "මම උදව් කරන්නම්."

- English message = English reply.
- Tamil message = Tamil reply.
- Never reply Singlish when the user is speaking Sinhala/Singlish.
- Convert Singlish Sinhala meaning into Sinhala Unicode.
- Keep Sinhala replies natural and readable.
- Do not mention this rule.

REPLY STYLE:
- Reply like ChatGPT.
- Keep answers concise.
- Do not write long essays unless user asks.
- Use simple explanations.
- Use bullet points when useful.
- Avoid unnecessary introductions.


You can help with:
- Coding
- Science
- Maths
- Technology
- General knowledge
- Writing
- Explanations
- Writing
- speech
- Exam preparations
- Learning new topics
- Language translations
- Learning new languages
- Learning new skills
- Life advice
- Learning new hobbies
- Learning new things
- Learning new subjects
- Learning new technologies
- Learning new programming languages
- Learning new frameworks
- Learning new tools
- Basic to advanced level topics
- Creating content
- Debugging code
- Giving step by step solutions
- Giving detailed explanations
- Giving concise answers
- Giving simple answers
- Giving easy to understand answers
- Giving in depth answers
- Giving short answers
- Giving long answers
- Giving examples
- Giving code examples
- Giving real life examples
- Giving practical examples
- Giving theoretical examples
- Giving mathematical examples
- Giving scientific examples
- Giving historical examples
- Giving philosophical examples
- Giving detailed explanations with examples
- Giving concise explanations with examples
- Giving simple explanations with examples
- Genarating new images based on user prompts
- Genarating new text based on user prompts
- Genarating new code based on user prompts
- Genarating new content based on user prompts
- Gebarate image captions based on user prompts
- Genarating new ideas based on user prompts
- Genarating new concepts based on user prompts
- Genarating new solutions based on user prompts
- and much more.
"""
            }
        ]

        for item in history:
            role = item.get("role")
            content = item.get("content") or ""
            if role in {"user", "assistant", "system"} and content:
                messages_payload.append({"role": role, "content": content})

        # Fetch both live web search and real-time data
        search_context = _build_live_search_context(user_message)
        real_time_context = _build_real_time_context(user_message)
        
        if attachment_context:
            user_prompt = f"{user_message}\n\n[Uploaded file context]\n{attachment_context}"
        else:
            user_prompt = user_message

        # Combine contexts with appropriate instructions
        context_parts = []
        if real_time_context:
            context_parts.append(real_time_context)
            context_parts.append("Use this real-time data to provide the most current and accurate information.")
        if search_context:
            context_parts.append(search_context)
            context_parts.append("If the search results do not clearly answer the question, state that you could not verify the latest information.")
        
        if context_parts:
            user_prompt = f"{user_prompt}\n\n" + "\n\n".join(context_parts)

        messages_payload.append({"role": "user", "content": user_prompt})

        gemini_client = _get_gemini_client()
        print("Route gemini client available:", gemini_client is not None)
        if gemini_client is None:
            return jsonify({
                "reply": "The AI service is currently unavailable because the Gemini API key is not configured on the server. Please contact the administrator to set GEMINI_API_KEY in the deployment environment."
            })

        model_name = _get_gemini_model()
        contents = _build_gemini_contents(messages_payload)

        print("=== Gemini request ===")
        print("Model:", model_name)
        print("Message count:", len(messages_payload))
        print("Environment:", os.getcwd())

        try:
            response = _generate_gemini_content(
                gemini_client,
                contents,
                model_name=model_name,
                config={
                    "temperature": 0.7,
                    "max_output_tokens": 800,
                    "response_mime_type": "text/plain",
                },
            )
            print("=== Gemini response received ===")
            print("Response type:", type(response))
            print("Response dir:", [x for x in dir(response) if not x.startswith('_')][:10])
            
            # Strategy 1: Direct text attribute
            answer = None
            if hasattr(response, 'text') and response.text:
                answer = response.text
                print("Strategy 1 (text attr): SUCCESS", len(answer), "chars")
            
            # Strategy 2: candidates array
            if not answer and hasattr(response, 'candidates') and response.candidates:
                print("Trying strategy 2 (candidates)...")
                first_candidate = response.candidates[0]
                print("First candidate type:", type(first_candidate))
                if hasattr(first_candidate, 'content') and first_candidate.content:
                    print("First candidate has content")
                    if hasattr(first_candidate.content, 'parts') and first_candidate.content.parts:
                        parts = first_candidate.content.parts
                        print(f"Parts count: {len(parts)}")
                        if parts and hasattr(parts[0], 'text'):
                            answer = parts[0].text
                            if answer:
                                print("Strategy 2 (candidates): SUCCESS", len(answer), "chars")
            
            # Strategy 3: dict conversion
            if not answer:
                print("Trying strategy 3 (dict conversion)...")
                try:
                    response_dict = dict(response) if hasattr(response, '__iter__') else None
                    if response_dict and 'text' in response_dict:
                        answer = response_dict['text']
                        if answer:
                            print("Strategy 3 (dict): SUCCESS", len(answer), "chars")
                except Exception as dict_err:
                    print("Strategy 3 failed:", dict_err)
            
            # Strategy 4: JSON serialization fallback
            if not answer:
                print("Trying strategy 4 (JSON fallback)...")
                try:
                    import json
                    response_json = json.dumps(response, default=str)
                    response_obj = json.loads(response_json)
                    if isinstance(response_obj, dict) and 'text' in response_obj:
                        answer = response_obj['text']
                        if answer:
                            print("Strategy 4 (JSON): SUCCESS", len(answer), "chars")
                except Exception as json_err:
                    print("Strategy 4 failed:", json_err)
            
            if not answer:
                print("⚠ WARNING: Could not extract text from response")
                print("Full response:", response)
                print("Response string:", str(response))
                answer = str(response) if response else "No response received from Gemini."
            
            print("=== Final answer ===")
            print("Length:", len(answer))
            preview_text = answer[:100] if len(answer) > 100 else answer
            print("Preview:", _safe_console_text(repr(preview_text)))
            print("=== Gemini full response ===")
            try:
                print(response)
            except Exception:
                pass

            # If the answer text contains known API error indicators, raise to expose full traceback
            lower_ans = (answer or "").lower()
            if any(keyword in lower_ans for keyword in ["quota", "insufficient_quota", "quota exceeded", "you exceeded your current quota", "rate_limit", "resource_exhausted"]):
                print("!!! Detected API error-like response in answer, raising to return full traceback")
                raise RuntimeError(f"Gemini API returned error-like response: {answer}")
        except Exception as chat_exc:
            print("=== Gemini request exception ===")
            print("Exception type:", type(chat_exc).__name__)
            print("Exception message:", _safe_console_text(chat_exc))
            traceback.print_exc()
            print("status_code:", getattr(chat_exc, "status_code", None))
            print("body:", getattr(chat_exc, "body", None))
            raise

        final_json = {"reply": answer}
        print("=== Final JSON to return ===")
        print(_safe_console_text(final_json))
        return jsonify(final_json)

    except Exception as e:
        print("=== Fatal error in chat route ===")
        print("Exception type:", type(e).__name__)
        print("Exception message:", _safe_console_text(e))
        tb_text = traceback.format_exc()
        print(_safe_console_text(tb_text))
        error_message = _format_gemini_error_message(e)
        print("Formatted error:", _safe_console_text(error_message))
        return jsonify({
            "reply": error_message,
            "error": str(e),
            "traceback": tb_text,
            "raw_error": repr(e)
        })


@app.route("/api/chat", methods=["POST"])
def api_chat():
    return chat()


if __name__=="__main__":
    debug_gemini_direct()
    app.run(
        host="0.0.0.0",
        port=int(os.getenv("PORT", "5001")),
        debug=True
    )